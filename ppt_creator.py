
##Library Files Include##
from pptx import Presentation
from gtts import gTTS
import win32com.client
##Notepad Inclusion ##
file1 = open("1.txt","r+") 
myText = file1.read()#"Jesus is my Lord"
#print(myText)



##Split the Note Pad
from text_split import split_into_sentences
a=split_into_sentences(myText);
prs = Presentation()
first_slide_layout = prs.slide_layouts[1]
title="THis is a Creation of GWU"
language = 'en'
## PPT Creation ##
# Create a ppt2jpg function
ppt_path=r"C:\Users\Immanual mech\Desktop\Thermal\Voice ppt"
def ppt2jpg():
    output_path = ppt_path #output path is the same as ppt path
    ppt_app = win32com.client.Dispatch('PowerPoint.Application')
    ppt = ppt_app.Presentations.Open(r'C:\Users\Immanual mech\Desktop\Thermal\Voice ppt\Output.pptx')  #start PowerPoint
    ppt.SaveAs(output_path, 17)  #17 is the number to save as jpg file type
    ppt_app.Quit()  # close PowerPoint

    
def add_slide(prs, first_slide_layout, title,content):
    """Return slide newly added to `prs` using `layout` and having `title`."""
    slide = prs.slides.add_slide(first_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    return slide

for i in range(len(a)):
    content = a[i];
    output= gTTS (text=content,lang = language, slow = False)
    output.save(str(i)+".mp3")
    add_slide(prs, first_slide_layout, title,content)
    prs.save("Output.pptx")


# Saving file

ppt2jpg() 
print("done")




